"""
Agrega 4 slides de arquitecturas de modelos al TESIS 22.pptx
respetando la paleta: #A3C5ED / #082836 / #F2F2F2
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

PATH = '/home/omar/Desktop/TESIS1/adversarial-attack-cafa/presentacion/TESIS 22.pptx'
prs  = Presentation(PATH)

# ── Paleta original del TESIS ────────────────────────────────────────────────
BLUE   = RGBColor(0xA3, 0xC5, 0xED)   # #A3C5ED  azul claro (header/accent)
DARK   = RGBColor(0x08, 0x28, 0x36)   # #082836  azul oscuro (bordes/texto)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)   # #F2F2F2  gris claro (fondo panel)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DGRAY  = RGBColor(0xD0, 0xD8, 0xE0)   # separadores
MGRAY  = RGBColor(0x60, 0x80, 0x96)   # texto secundario

BLANK = prs.slide_layouts[6]

W = prs.slide_width.inches    # 13.33"
H = prs.slide_height.inches   # 7.50"

# ── Helpers ──────────────────────────────────────────────────────────────────
def bg(sl):
    s = sl.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.fill.background()

def rect(sl, x, y, w, h, fill=LGRAY, border=DARK, lw=1.0):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(lw)
    return s

def rnd(sl, x, y, w, h, fill=LGRAY, border=DARK, lw=1.2):
    s = sl.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(lw)
    return s

def txt(sl, text, x, y, w, h, size=10, bold=False, color=DARK,
        align=PP_ALIGN.CENTER, italic=False):
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return txb

def hdr(sl, title, subtitle=''):
    """Encabezado azul idéntico al del slide original."""
    rect(sl, 0, 0, W, 0.72, fill=DARK, border=DARK, lw=0)
    txt(sl, title, 0.25, 0.10, 10, 0.52, size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(sl, subtitle, 0.25, 0.42, 12, 0.28, size=9, color=BLUE, align=PP_ALIGN.LEFT, italic=True)
    # barra de acento inferior
    rect(sl, 0, 0.68, W, 0.05, fill=BLUE, border=BLUE, lw=0)

def block(sl, label, x, y, w=2.0, h=0.45, fill=BLUE, bdr=DARK):
    """Bloque con header azul (estilo original)."""
    rnd(sl, x, y, w, h, fill=fill, border=bdr, lw=1.2)
    txt(sl, label, x, y+0.04, w, h-0.08, size=9, bold=True, color=DARK)

def block_gray(sl, label, x, y, w=2.0, h=0.45):
    """Bloque relleno gris (estilo caja original)."""
    rnd(sl, x, y, w, h, fill=LGRAY, border=DARK, lw=1.0)
    txt(sl, label, x, y+0.04, w, h-0.08, size=9, color=DARK)

def arrow_v(sl, cx, y1, y2):
    c = sl.shapes.add_connector(1,
        Inches(cx), Inches(y1), Inches(cx), Inches(y2))
    c.line.color.rgb = DARK; c.line.width = Pt(1.2)

def arrow_h(sl, x1, x2, cy):
    c = sl.shapes.add_connector(1,
        Inches(x1), Inches(cy), Inches(x2), Inches(cy))
    c.line.color.rgb = DARK; c.line.width = Pt(1.2)

def pill(sl, label, x, y, w=1.3, h=0.32, fill=LGRAY):
    rnd(sl, x, y, w, h, fill=fill, border=DARK, lw=0.8)
    txt(sl, label, x, y+0.02, w, h-0.04, size=8, color=DARK)

def spec_row(sl, k, v, x, y, wk=1.9, wv=2.5):
    txt(sl, k, x,    y, wk, 0.32, size=8.5, bold=True,  color=MGRAY, align=PP_ALIGN.LEFT)
    txt(sl, v, x+wk, y, wv, 0.32, size=8.5, bold=False, color=DARK,  align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE A — Regresión Logística
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
hdr(sl, "Regresión Logística", "Clasificador lineal · una capa · diferenciable (compatible con CaFA, HSJ, Boundary, Square)")

# Panel diagrama (izquierda)
rect(sl, 0.2, 0.82, 8.2, 6.5, fill=LGRAY, border=DARK, lw=1.0)
txt(sl, "ARQUITECTURA", 0.35, 0.85, 3, 0.28, size=8, bold=True, color=DARK, align=PP_ALIGN.LEFT)

# Nodos entrada
inputs = ["x₁","x₂","x₃","·","xₙ"]
iy_list = [1.5, 2.1, 2.7, 3.3, 3.9]
for lbl, iy in zip(inputs, iy_list):
    s = sl.shapes.add_shape(9, Inches(1.0), Inches(iy-0.22),
                             Inches(0.44), Inches(0.44))
    s.fill.solid(); s.fill.fore_color.rgb = BLUE
    s.line.color.rgb = DARK; s.line.width = Pt(1.0)
    txt(sl, lbl, 1.0, iy-0.20, 0.44, 0.40, size=10, bold=True, color=DARK)

# Pesos
for i, iy in enumerate(iy_list):
    arrow_h(sl, 1.44, 3.9, iy)
    txt(sl, f"w{i+1 if i<4 else 'ₙ'}", 2.1+i*0.18, iy-0.22,
        0.45, 0.22, size=7.5, color=MGRAY, italic=True)

# Nodo suma
s = sl.shapes.add_shape(9, Inches(3.9), Inches(2.5),
                         Inches(0.7), Inches(0.7))
s.fill.solid(); s.fill.fore_color.rgb = BLUE
s.line.color.rgb = DARK; s.line.width = Pt(1.5)
txt(sl, "Σ", 3.9, 2.52, 0.70, 0.66, size=20, bold=True, color=DARK)
txt(sl, "Σwᵢxᵢ + b", 3.7, 3.28, 1.2, 0.25, size=7.5, color=MGRAY, italic=True)

# flecha a activación
arrow_h(sl, 4.60, 5.5, 2.85)
txt(sl, "logits", 4.65, 2.6, 0.9, 0.22, size=7.5, color=MGRAY, italic=True)

# Nodo Softmax
s = sl.shapes.add_shape(9, Inches(5.5), Inches(2.5),
                         Inches(0.7), Inches(0.7))
s.fill.solid(); s.fill.fore_color.rgb = BLUE
s.line.color.rgb = DARK; s.line.width = Pt(1.5)
txt(sl, "σ", 5.5, 2.52, 0.70, 0.66, size=20, bold=True, color=DARK)
txt(sl, "Softmax", 5.42, 3.28, 0.9, 0.25, size=7.5, color=MGRAY, italic=True)

# Salidas
arrow_h(sl, 6.20, 7.2, 2.85)
for oy, ol in [(2.2, "P(legítimo)"), (3.5, "P(fraude)")]:
    c = sl.shapes.add_connector(1, Inches(7.2), Inches(2.85), Inches(7.2), Inches(oy+0.15))
    c.line.color.rgb = DARK; c.line.width = Pt(1.0)
    s = sl.shapes.add_shape(9, Inches(7.2), Inches(oy),
                             Inches(0.44), Inches(0.44))
    s.fill.solid(); s.fill.fore_color.rgb = LGRAY
    s.line.color.rgb = DARK; s.line.width = Pt(1.0)
    txt(sl, ol, 7.7, oy, 1.2, 0.44, size=8.5, color=DARK, align=PP_ALIGN.LEFT)

# Fórmula
rect(sl, 0.4, 4.6, 7.8, 0.6, fill=WHITE, border=BLUE, lw=1.2)
txt(sl, "f(x) = Wx + b    →    ŷ = Softmax(f(x))", 0.5, 4.65, 7.5, 0.5,
    size=11, bold=True, color=DARK)

# Panel specs (derecha)
rect(sl, 8.55, 0.82, 4.6, 6.5, fill=WHITE, border=DARK, lw=1.0)
txt(sl, "ESPECIFICACIONES", 8.7, 0.85, 4.3, 0.28, size=8, bold=True, color=DARK, align=PP_ALIGN.LEFT)
block(sl, "Arquitectura", 8.7, 1.2, 4.2, 0.38)
specs_lr = [
    ("Capas",         "1 × Linear(input_dim → 2)"),
    ("Activación",    "Softmax (implícita en CrossEntropy)"),
    ("Parámetros",    "input_dim × 2  +  2  (bias)"),
    ("",              ""),
    ("Optimizador",   "Adam"),
    ("lr",            "HPO: [1×10⁻⁵,  1×10⁻¹]"),
    ("weight_decay",  "HPO: [1×10⁻⁶,  1×10⁻²]"),
    ("Épocas máx.",   "100"),
    ("Batch size",    "4 096"),
    ("HPO trials",    "50  (Optuna)"),
    ("Monitor",       "AUC-ROC (validación)"),
    ("",              ""),
    ("Diferenciable", "✓  Sí  →  white-box + black-box"),
    ("Ataques",       "CaFA · HSJ · Boundary · Square"),
]
for i, (k, v) in enumerate(specs_lr):
    spec_row(sl, k, v, 8.7, 1.65 + i*0.36)

# Nota
rect(sl, 8.6, 6.55, 4.45, 0.65, fill=RGBColor(0xE8,0xF4,0xFF), border=BLUE, lw=1.0)
txt(sl, "Modelo más interpretable. Sirve como línea base para comparar la robustez de los modelos más complejos.",
    8.7, 6.58, 4.3, 0.60, size=8, color=DARK, italic=True, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE B — MLP
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
hdr(sl, "MLP — Perceptrón Multicapa", "Red neuronal profunda · BatchNorm + ReLU · diferenciable")

rect(sl, 0.2, 0.82, 8.2, 6.5, fill=LGRAY, border=DARK, lw=1.0)
txt(sl, "ARQUITECTURA  (n_layers = 3, hidden_dim = 128)", 0.35, 0.85, 7.8, 0.28,
    size=8, bold=True, color=DARK, align=PP_ALIGN.LEFT)

layers = [
    ("INPUT  (input_dim)",         LGRAY),
    ("BatchNorm1d (input_dim)",    BLUE),
    ("Linear  input_dim → 128",    BLUE),
    ("ReLU",                       BLUE),
    ("Linear  128 → 128",          BLUE),
    ("ReLU",                       BLUE),
    ("Linear  128 → 2  (logits)",  BLUE),
    ("Softmax  →  OUTPUT",         LGRAY),
]
bx, bw, bh = 3.0, 2.2, 0.52
gap = 0.18
total_h = len(layers) * bh + (len(layers)-1) * gap
start_y = (H - total_h) / 2 + 0.3

for i, (lbl, fill) in enumerate(layers):
    by = start_y + i * (bh + gap)
    rnd(sl, bx, by, bw, bh, fill=fill, border=DARK, lw=1.2)
    txt(sl, lbl, bx, by+0.06, bw, bh-0.12, size=9, bold=(fill==BLUE), color=DARK)
    if i < len(layers) - 1:
        arrow_v(sl, bx + bw/2, by + bh, by + bh + gap)

# Corchete de repetición
rep_y1 = start_y + 4 * (bh + gap)
rep_y2 = start_y + 5 * (bh + gap) + bh
rx = bx + bw + 0.15
rect(sl, rx, rep_y1, 2.5, rep_y2 - rep_y1, fill=RGBColor(0xD8,0xEA,0xF8), border=BLUE, lw=1.0)
txt(sl, "× (n_layers − 1)\nrepeticiones\n(HPO: 2 a 5 capas)", rx+0.08, rep_y1+0.05, 2.35, rep_y2-rep_y1-0.1,
    size=8, color=DARK, italic=True)

# Panel specs
rect(sl, 8.55, 0.82, 4.6, 6.5, fill=WHITE, border=DARK, lw=1.0)
txt(sl, "ESPECIFICACIONES", 8.7, 0.85, 4.3, 0.28, size=8, bold=True, color=DARK, align=PP_ALIGN.LEFT)
block(sl, "Arquitectura", 8.7, 1.2, 4.2, 0.38)
specs_mlp = [
    ("Normaliz.",    "BatchNorm1d en capa de entrada"),
    ("Capas ocultas","n_layers  ∈  [2, 5]  (HPO)"),
    ("Neuronas/capa","hidden_dim  ∈  [32, 512]  (HPO)"),
    ("Activación",   "ReLU (ocultas) · Softmax (salida)"),
    ("",             ""),
    ("Pérdida",      "CrossEntropy ponderada (class_weight)"),
    ("Optimizador",  "Adam"),
    ("lr",           "HPO: [1×10⁻⁷,  1×10⁻¹]"),
    ("weight_decay", "HPO: [1×10⁻⁶,  1×10⁻¹]"),
    ("Épocas máx.",  "100"),
    ("Batch size",   "2 048"),
    ("HPO trials",   "100  (Optuna)"),
    ("Monitor",      "AUC-ROC (validación)"),
    ("",             ""),
    ("Diferenciable","✓  Sí  →  white-box + black-box"),
]
for i, (k, v) in enumerate(specs_mlp):
    spec_row(sl, k, v, 8.7, 1.65 + i*0.33)

rect(sl, 8.6, 6.6, 4.45, 0.60, fill=RGBColor(0xE8,0xF4,0xFF), border=BLUE, lw=1.0)
txt(sl, "La BatchNorm en entrada estabiliza el entrenamiento con datos tabulares de distinta escala.",
    8.7, 6.63, 4.3, 0.55, size=8, color=DARK, italic=True, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE C — LSTM-Attention
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
hdr(sl, "LSTM con Mecanismo de Atención (LSTM-Att.)", "Red recurrente · atención aditiva · diferenciable")

rect(sl, 0.2, 0.82, 8.2, 6.5, fill=LGRAY, border=DARK, lw=1.0)
txt(sl, "ARQUITECTURA  (hidden_dim = 64, n_lstm_layers = 2, dropout = 0.3)", 0.35, 0.85, 7.8, 0.28,
    size=8, bold=True, color=DARK, align=PP_ALIGN.LEFT)

# Flujo principal
flow = [
    ("INPUT  (batch, input_dim)",               LGRAY, 0.50),
    ("Reshape  →  (batch, seq_len, step_dim)",  BLUE,  0.48),
    ("LSTM × 2 capas\nhidden_dim = 64\ndropout = 0.3\n→  (batch, seq, 64)", BLUE, 0.92),
    ("Atención Aditiva\ntanh → Linear(64,1) → Softmax\n→  Σαₜhₜ  =  contexto (batch,64)", BLUE, 0.88),
    ("Dropout (0.3)  +  Linear (64 → 2)",       BLUE,  0.48),
    ("OUTPUT  logits  →  Softmax",              LGRAY, 0.48),
]
bx2, bw2 = 2.6, 2.8
gap2 = 0.14
cy = 1.1
for i, (lbl, fill, bh2) in enumerate(flow):
    rnd(sl, bx2, cy, bw2, bh2, fill=fill, border=DARK, lw=1.2)
    txt(sl, lbl, bx2, cy+0.04, bw2, bh2-0.08, size=8.5, bold=(fill==BLUE), color=DARK)
    next_y = cy + bh2
    if i < len(flow)-1:
        arrow_v(sl, bx2+bw2/2, next_y, next_y+gap2)
    cy = next_y + gap2

# Detalle fórmulas atención (recuadro derecho del panel)
rect(sl, 5.7, 2.85, 2.55, 2.0, fill=WHITE, border=BLUE, lw=1.2)
txt(sl, "Atención Aditiva:", 5.8, 2.88, 2.35, 0.28, size=8, bold=True, color=DARK, align=PP_ALIGN.LEFT)
formulas = [
    "eₜ = tanh(W·hₜ + b)",
    "αₜ = softmax(e)ₜ",
    "c  = Σ αₜ · hₜ",
]
for j, f in enumerate(formulas):
    txt(sl, f, 5.82, 3.22+j*0.46, 2.3, 0.42, size=9, color=DARK, align=PP_ALIGN.LEFT, italic=True)

# Panel specs
rect(sl, 8.55, 0.82, 4.6, 6.5, fill=WHITE, border=DARK, lw=1.0)
txt(sl, "ESPECIFICACIONES", 8.7, 0.85, 4.3, 0.28, size=8, bold=True, color=DARK, align=PP_ALIGN.LEFT)
block(sl, "Arquitectura", 8.7, 1.2, 4.2, 0.38)
specs_lstm = [
    ("Capas LSTM",    "n_lstm_layers  ∈  [1, 4]  (HPO)"),
    ("Unidades",      "hidden_dim  ∈  [32, 256]  (HPO)"),
    ("Dropout",       "∈  [0.0,  0.5]  (HPO)"),
    ("Seq. length",   "1  (transacción como secuencia)"),
    ("Atención",      "Aditiva  — tanh + Softmax"),
    ("",              ""),
    ("Pérdida",       "CrossEntropy ponderada (class_weight)"),
    ("Optimizador",   "Adam"),
    ("lr",            "HPO: [1×10⁻⁵,  1×10⁻²]"),
    ("weight_decay",  "HPO: [1×10⁻⁶,  1×10⁻²]"),
    ("Épocas máx.",   "100"),
    ("Batch size",    "1 024"),
    ("HPO trials",    "50  (Optuna)"),
    ("Monitor",       "AUC-ROC (validación)"),
    ("",              ""),
    ("Diferenciable", "✓  Sí  →  white-box + black-box"),
]
for i, (k, v) in enumerate(specs_lstm):
    spec_row(sl, k, v, 8.7, 1.65 + i*0.31)

rect(sl, 8.6, 6.6, 4.45, 0.60, fill=RGBColor(0xE8,0xF4,0xFF), border=BLUE, lw=1.0)
txt(sl, "El mecanismo de atención pondera la importancia de cada paso temporal del LSTM antes de clasificar.",
    8.7, 6.63, 4.3, 0.55, size=8, color=DARK, italic=True, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE D — XGBoost
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
hdr(sl, "XGBoost — Gradient Boosting de Árboles", "Ensamble de árboles · NO diferenciable · solo ataques black-box (HSJ, Boundary, Square)")

rect(sl, 0.2, 0.82, 8.2, 6.5, fill=LGRAY, border=DARK, lw=1.0)
txt(sl, "ARQUITECTURA  (n_estimators = 100, max_depth = 6)", 0.35, 0.85, 7.8, 0.28,
    size=8, bold=True, color=DARK, align=PP_ALIGN.LEFT)

# INPUT
rnd(sl, 2.8, 1.1, 2.6, 0.50, fill=LGRAY, border=DARK, lw=1.2)
txt(sl, "INPUT  vector de features x", 2.8, 1.13, 2.6, 0.44, size=9, color=DARK)
arrow_v(sl, 4.1, 1.60, 1.95)

# 3 árboles
tree_info = [("Árbol 1", 1.1), ("Árbol 2", 3.1), ("Árbol K\n(K=100)", 5.9)]
for lbl, tx in tree_info:
    rnd(sl, tx, 1.95, 1.7, 0.65, fill=BLUE, border=DARK, lw=1.2)
    txt(sl, lbl, tx, 1.98, 1.7, 0.60, size=9, bold=True, color=DARK)
    c = sl.shapes.add_connector(1, Inches(4.1), Inches(1.60), Inches(tx+0.85), Inches(1.95))
    c.line.color.rgb = DARK; c.line.width = Pt(1.0)

txt(sl, "···", 4.0, 2.1, 0.6, 0.4, size=16, bold=True, color=MGRAY)

# Suma
rnd(sl, 2.85, 3.2, 2.5, 0.55, fill=BLUE, border=DARK, lw=1.5)
txt(sl, "Σ predicciones ponderadas", 2.85, 3.23, 2.5, 0.48, size=9, bold=True, color=DARK)
for tx in [1.95, 4.0, 6.75]:
    c = sl.shapes.add_connector(1, Inches(tx), Inches(2.60), Inches(4.1), Inches(3.20))
    c.line.color.rgb = DARK; c.line.width = Pt(1.0)
arrow_v(sl, 4.1, 3.75, 4.15)

# Sigmoid
rnd(sl, 2.85, 4.15, 2.5, 0.52, fill=BLUE, border=DARK, lw=1.2)
txt(sl, "Sigmoid  →  P(fraude)", 2.85, 4.18, 2.5, 0.45, size=9, bold=True, color=DARK)
arrow_v(sl, 4.1, 4.67, 5.05)

# Output
rnd(sl, 2.85, 5.05, 2.5, 0.50, fill=LGRAY, border=DARK, lw=1.2)
txt(sl, "OUTPUT: P(legítimo) · P(fraude)", 2.85, 5.08, 2.5, 0.44, size=9, color=DARK)

# Detalle árbol
rect(sl, 0.35, 3.9, 2.35, 2.8, fill=WHITE, border=BLUE, lw=1.2)
txt(sl, "ÁRBOL DE DECISIÓN:", 0.45, 3.93, 2.1, 0.28, size=7.5, bold=True, color=DARK, align=PP_ALIGN.LEFT)
txt(sl, "if  x[j] > θ  →  rama derecha\nelse           →  rama izquierda\n\nCada hoja devuelve un\nvalor de predicción.\nmax_depth = 6\n→ hasta 64 hojas por árbol", 0.45, 4.25, 2.15, 2.35,
    size=8, color=DARK, align=PP_ALIGN.LEFT, italic=False)

# Alerta no diferenciable
rect(sl, 5.65, 3.9, 2.65, 2.8, fill=RGBColor(0xFF,0xF0,0xF0), border=RGBColor(0xCC,0x33,0x33), lw=1.5)
txt(sl, "⚠  NO DIFERENCIABLE", 5.75, 3.93, 2.45, 0.30, size=8.5, bold=True,
    color=RGBColor(0xAA,0x22,0x22), align=PP_ALIGN.LEFT)
txt(sl,
    "Las decisiones x[j] > θ son\n"
    "funciones escalón discontinuas.\n"
    "No existe ∂f/∂x.\n\n"
    "→ CaFA incompatible.\n"
    "→ HSJ, Boundary y Square\n"
    "   Attack SÍ compatibles\n"
    "   (solo consultan predicción).",
    5.75, 4.28, 2.5, 2.35, size=8, color=RGBColor(0x77,0x11,0x11), align=PP_ALIGN.LEFT)

# Panel specs
rect(sl, 8.55, 0.82, 4.6, 6.5, fill=WHITE, border=DARK, lw=1.0)
txt(sl, "ESPECIFICACIONES", 8.7, 0.85, 4.3, 0.28, size=8, bold=True, color=DARK, align=PP_ALIGN.LEFT)
block(sl, "Arquitectura", 8.7, 1.2, 4.2, 0.38)
specs_xgb = [
    ("Tipo",           "XGBClassifier  (xgboost)"),
    ("Estimadores",    "n_estimators  ∈  [100, 600]  (HPO)"),
    ("Profundidad",    "max_depth  ∈  [3, 10]  (HPO)"),
    ("Learning rate",  "∈  [0.01,  0.3]  (HPO)"),
    ("min_child_w.",   "∈  [1, 20]  (HPO)"),
    ("subsample",      "∈  [0.5, 1.0]  (HPO)"),
    ("colsample",      "∈  [0.5, 1.0]  (HPO)"),
    ("",               ""),
    ("Objetivo",       "binary:logistic"),
    ("Métrica eval.",  "AUC-PR  (aucpr)"),
    ("Desbalance",     "scale_pos_weight = peso minoritaria"),
    ("",               ""),
    ("HPO trials",     "30  (Optuna TPE, seed=42)"),
    ("",               ""),
    ("Diferenciable",  "✗  No  →  solo black-box"),
    ("Ataques",        "HSJ · Boundary · Square Attack"),
]
for i, (k, v) in enumerate(specs_xgb):
    spec_row(sl, k, v, 8.7, 1.65 + i*0.31)

rect(sl, 8.6, 6.6, 4.45, 0.60, fill=RGBColor(0xFF,0xF0,0xF0), border=RGBColor(0xCC,0x33,0x33), lw=1.0)
txt(sl, "XGBoost usa gradientes para ENTRENAR (ajustar árboles), pero no los expone sobre el input. Su función de predicción es discontinua.",
    8.7, 6.63, 4.3, 0.55, size=8, color=RGBColor(0x88,0x22,0x22), italic=True, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# Guardar
# ════════════════════════════════════════════════════════════════════════════
prs.save(PATH)
print(f"Guardado: {PATH}  ({len(prs.slides)} slides)")
