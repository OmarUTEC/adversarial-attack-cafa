"""
Genera arquitecturas_modelos.pptx con diagramas de los 4 modelos de clasificación.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import math

# ── Paleta ──────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x17, 0x2A)   # fondo oscuro
CARD      = RGBColor(0x1E, 0x29, 0x3B)   # caja
BORDER    = RGBColor(0x33, 0x4E, 0x6B)   # borde
TITLE_CLR = RGBColor(0xF8, 0xFA, 0xFC)   # blanco
SUB_CLR   = RGBColor(0x94, 0xA3, 0xB8)   # gris claro
ACC_BLUE  = RGBColor(0x38, 0xBD, 0xF8)   # acento azul
ACC_GREEN = RGBColor(0x34, 0xD3, 0x99)   # acento verde
ACC_PURP  = RGBColor(0xA7, 0x8B, 0xFA)   # acento morado
ACC_ORG   = RGBColor(0xFB, 0xBF, 0x24)   # acento naranja
RED_CLR   = RGBColor(0xF8, 0x71, 0x85)   # rojo
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# Colores por modelo
MODEL_COLORS = {
    'lr':    ACC_BLUE,
    'mlp':   ACC_GREEN,
    'lstm':  ACC_PURP,
    'xgb':   ACC_ORG,
}

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completamente en blanco


# ════════════════════════════════════════════════════════════════════════════
# Helpers
# ════════════════════════════════════════════════════════════════════════════

def bg(slide):
    """Fondo oscuro para el slide."""
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    return shape


def box(slide, x, y, w, h, fill=CARD, border=BORDER, radius=0.12):
    """Caja redondeada."""
    s = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(1.2)
    return s


def label(slide, text, x, y, w, h,
          size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER, italic=False):
    """Cuadro de texto."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def circle(slide, cx, cy, r, fill, border=WHITE):
    """Círculo centrado en (cx, cy) con radio r (pulgadas)."""
    s = slide.shapes.add_shape(9,
        Inches(cx - r), Inches(cy - r), Inches(r*2), Inches(r*2))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = border; s.line.width = Pt(1.2)
    return s


def arrow(slide, x1, y1, x2, y2, color=SUB_CLR):
    """Flecha de (x1,y1) a (x2,y2) en pulgadas."""
    from pptx.util import Pt
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    return connector


def slide_title(slide, title, subtitle, accent):
    """Encabezado del slide."""
    # barra de acento izquierda
    bar = slide.shapes.add_shape(1, Inches(0.3), Inches(0.22),
                                  Inches(0.06), Inches(0.56))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    label(slide, title,    0.44, 0.18, 12, 0.38,
          size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    label(slide, subtitle, 0.44, 0.54, 12, 0.28,
          size=11, color=SUB_CLR, align=PP_ALIGN.LEFT, italic=True)


def footer(slide):
    label(slide,
          "Modelos de Clasificación para Detección de Fraude Financiero  ·  Framework Adversarial CaFA",
          0.2, 7.18, 12.9, 0.25, size=8, color=SUB_CLR)


# ════════════════════════════════════════════════════════════════════════════
# Slide 1 — Portada
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)

# Gradiente simulado con rectángulos
for i, alpha in enumerate([0.04, 0.03, 0.02]):
    s = sl.shapes.add_shape(1, Inches(0), Inches(i*2.5),
                             prs.slide_width, Inches(2.5))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x1E, 0x29+i*8, 0x3B+i*8)
    s.line.fill.background()

# Título principal
label(sl, "Arquitecturas de Modelos de Clasificación",
      0.5, 1.8, 12.3, 0.8, size=34, bold=True, color=WHITE)
label(sl, "Framework de Evaluación de Robustez Adversarial en Detección de Fraude Financiero",
      0.5, 2.65, 12.3, 0.45, size=16, color=SUB_CLR, italic=True)

# Línea separadora
sep = sl.shapes.add_shape(1, Inches(0.5), Inches(3.25), Inches(12.3), Inches(0.03))
sep.fill.solid(); sep.fill.fore_color.rgb = ACC_BLUE
sep.line.fill.background()

# Chips de modelos
chips = [("Logistic Regression", ACC_BLUE), ("MLP", ACC_GREEN),
         ("LSTM-Attention", ACC_PURP), ("XGBoost", ACC_ORG)]
for i, (name, clr) in enumerate(chips):
    cx = 1.6 + i * 2.85
    s = sl.shapes.add_shape(5, Inches(cx), Inches(3.6), Inches(2.4), Inches(0.55))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(
        min(clr[0]+20, 255) if isinstance(clr, tuple) else 0x1E,
        0x29, 0x3B)
    # draw colored chip
    s2 = sl.shapes.add_shape(5, Inches(cx), Inches(3.6), Inches(2.4), Inches(0.55))
    s2.fill.solid(); s2.fill.fore_color.rgb = CARD
    s2.line.color.rgb = clr; s2.line.width = Pt(1.5)
    label(sl, name, cx+0.08, 3.67, 2.24, 0.38,
          size=12, bold=True, color=clr)

label(sl, "PyTorch Lightning  ·  XGBoost  ·  Optimización: Optuna  ·  Entrenamiento: hasta 100 épocas",
      0.5, 4.4, 12.3, 0.35, size=11, color=SUB_CLR, italic=True)

# Detalles técnicos
details = [
    ("Datasets", "Credit Card 2023 · AMLworld HI-Small"),
    ("Preprocesamiento", "One-Hot Encoding · BatchNorm · SMOTE (variante)"),
    ("Optimizador", "Adam  (lr variable según HPO)"),
    ("Criterio", "Cross-Entropy ponderada por clase  ·  AUC-ROC"),
]
for i, (k, v) in enumerate(details):
    row_y = 5.0 + i * 0.42
    label(sl, k + ":", 1.2, row_y, 2.2, 0.36, size=10, bold=True, color=ACC_BLUE, align=PP_ALIGN.RIGHT)
    label(sl, v,       3.5, row_y, 9.0, 0.36, size=10, color=WHITE, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# Slide 2 — Logistic Regression
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Regresión Logística", "Clasificador lineal diferenciable · baseline interpretable", ACC_BLUE)
footer(sl)

# Panel izquierdo: diagrama
box(sl, 0.25, 0.95, 7.8, 6.3, border=ACC_BLUE)
label(sl, "ARQUITECTURA", 0.35, 0.98, 3, 0.28, size=8, color=ACC_BLUE, bold=True, align=PP_ALIGN.LEFT)

# Nodos de entrada
n_inputs = 5
input_labels = ["x₁", "x₂", "x₃", "·", "xₙ"]
input_y = [1.55, 2.15, 2.75, 3.35, 3.95]
r = 0.22
for i, (lbl, iy) in enumerate(zip(input_labels, input_y)):
    circle(sl, 1.35, iy, r, CARD, ACC_BLUE)
    label(sl, lbl, 1.35-r, iy-0.18, r*2, 0.36, size=10, bold=True, color=ACC_BLUE)

# Pesos
for i, iy in enumerate(input_y):
    arrow(sl, 1.57, iy, 4.05, 2.82, ACC_BLUE)
    label(sl, f"w{i+1 if i<4 else 'ₙ'}", 2.5 + i*0.05, iy - 0.05 + i*0.12,
          0.35, 0.25, size=8, color=SUB_CLR)

# Nodo suma
circle(sl, 4.35, 2.82, 0.35, RGBColor(0x17, 0x3A, 0x5C), ACC_BLUE)
label(sl, "Σ", 4.35-0.35, 2.82-0.25, 0.70, 0.50, size=18, bold=True, color=ACC_BLUE)
label(sl, "Σwᵢxᵢ + b", 3.82, 3.35, 1.1, 0.28, size=8, color=SUB_CLR, italic=True)

# Flecha a activación
arrow(sl, 4.70, 2.82, 5.6, 2.82, ACC_BLUE)
label(sl, "logits", 4.75, 2.55, 0.8, 0.25, size=8, color=SUB_CLR, italic=True)

# Nodo Softmax
circle(sl, 5.95, 2.82, 0.35, RGBColor(0x0C, 0x40, 0x6A), ACC_BLUE)
label(sl, "σ", 5.95-0.35, 2.82-0.25, 0.70, 0.50, size=18, bold=True, color=WHITE)
label(sl, "Softmax", 5.62, 3.3, 0.72, 0.25, size=8, color=SUB_CLR, italic=True)

# Flecha a salida
arrow(sl, 6.30, 2.82, 7.0, 2.82, ACC_BLUE)

# Salidas
out_y = [2.22, 3.42]
out_labels = ["P(legítimo)", "P(fraude)"]
for oy, ol in zip(out_y, out_labels):
    arrow(sl, 7.0, 2.82, 7.15, oy, ACC_GREEN)
    circle(sl, 7.45, oy, 0.22, CARD, ACC_GREEN)
    label(sl, ol, 7.7, oy-0.15, 1.0, 0.30, size=9, color=ACC_GREEN, align=PP_ALIGN.LEFT)

# Etiqueta sesgo
label(sl, "+ b (bias)", 3.85, 4.4, 1.1, 0.28, size=9, color=SUB_CLR, italic=True)

# Panel derecho: especificaciones
box(sl, 8.3, 0.95, 4.75, 6.3, border=BORDER)
label(sl, "ESPECIFICACIONES", 8.45, 0.98, 4.5, 0.28, size=8, color=ACC_BLUE, bold=True, align=PP_ALIGN.LEFT)

specs = [
    ("Capa",             "Linear(input_dim → 2)"),
    ("Función",          "f(x) = Wx + b"),
    ("Activación salida","Softmax (implícita en\nCross-Entropy)"),
    ("Parámetros",       "input_dim × 2 + 2"),
    ("Pérdida",          "Cross-Entropy ponderada"),
    ("Optimizador",      "Adam"),
    ("lr",               "HPO: [1e-5, 1e-1]"),
    ("weight_decay",     "HPO: [1e-6, 1e-2]"),
    ("Épocas máx.",      "100"),
    ("Batch size",       "4 096"),
    ("Monitor HPO",      "AUC-ROC (val)"),
    ("Diferenciable",    "✓  Sí — compatible con\nataques white-box"),
]
for i, (k, v) in enumerate(specs):
    ry = 1.38 + i * 0.44
    label(sl, k, 8.45, ry, 1.85, 0.40, size=9, bold=True, color=SUB_CLR, align=PP_ALIGN.LEFT)
    label(sl, v, 10.3, ry, 2.65, 0.40, size=9, color=WHITE, align=PP_ALIGN.LEFT)

# Nota al pie del panel
note = ("La regresión logística es el clasificador más simple e interpretable. "
        "Al ser completamente diferenciable, todos los ataques adversariales "
        "(white-box y black-box) son aplicables.")
box(sl, 8.35, 6.08, 4.6, 1.1, fill=RGBColor(0x0C, 0x2A, 0x47), border=ACC_BLUE)
label(sl, note, 8.45, 6.1, 4.45, 1.05, size=8, color=SUB_CLR, italic=True, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# Slide 3 — MLP
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "MLP — Perceptrón Multicapa",
            "Red neuronal profunda con BatchNorm y activaciones ReLU", ACC_GREEN)
footer(sl)

box(sl, 0.25, 0.95, 7.8, 6.3, border=ACC_GREEN)
label(sl, "ARQUITECTURA  (n_layers=3, hidden_dim=128)", 0.35, 0.98, 7.5, 0.28,
      size=8, color=ACC_GREEN, bold=True, align=PP_ALIGN.LEFT)

# Capas del MLP como bloques verticales
layers_mlp = [
    ("INPUT\ninput_dim",  SUB_CLR,  0.75),
    ("BatchNorm1d\n(input_dim)", ACC_BLUE, 1.55),
    ("Linear\ninput→128",  ACC_GREEN, 2.5),
    ("ReLU",               ACC_GREEN, 3.25),
    ("Linear\n128→128",    ACC_GREEN, 4.0),
    ("ReLU",               ACC_GREEN, 4.75),
    ("Linear\n128→2",      ACC_GREEN, 5.55),
    ("OUTPUT\n(logits)",   SUB_CLR,  6.45),
]

bw, bh = 1.45, 0.55
bx = 3.3
for i, (lbl, clr, by) in enumerate(layers_mlp):
    s = sl.shapes.add_shape(5, Inches(bx), Inches(by), Inches(bw), Inches(bh))
    s.fill.solid(); s.fill.fore_color.rgb = CARD
    s.line.color.rgb = clr; s.line.width = Pt(1.5)
    label(sl, lbl, bx, by+0.03, bw, bh-0.06, size=9, bold=(clr != SUB_CLR), color=clr)
    if i < len(layers_mlp) - 1:
        arrow(sl, bx + bw/2, by + bh, bx + bw/2, layers_mlp[i+1][2], ACC_GREEN)

# Repetición indicada
bracket_x = bx + bw + 0.12
box(sl, bracket_x, 3.92, 2.2, 1.42, fill=RGBColor(0x0D, 0x2A, 0x1A), border=ACC_GREEN)
label(sl, "× (n_layers − 1)\n= 2 repeticiones\n(configurable 2–5)",
      bracket_x + 0.05, 4.0, 2.1, 1.26, size=8, color=ACC_GREEN, italic=True)

# Panel derecho: specs
box(sl, 8.3, 0.95, 4.75, 6.3, border=BORDER)
label(sl, "ESPECIFICACIONES", 8.45, 0.98, 4.5, 0.28, size=8, color=ACC_GREEN, bold=True, align=PP_ALIGN.LEFT)

specs = [
    ("Capas ocultas",     "n_layers (HPO: 2–5)"),
    ("Neuronas/capa",     "hidden_dim (HPO: 32–512)"),
    ("Normalización",     "BatchNorm1d en entrada"),
    ("Activación",        "ReLU (capas ocultas)"),
    ("Activación salida", "Softmax (en pérdida)"),
    ("Pérdida",           "Cross-Entropy ponderada"),
    ("Optimizador",       "Adam"),
    ("lr",                "HPO: [1e-7, 1e-1]"),
    ("weight_decay",      "HPO: [1e-6, 1e-1]"),
    ("Épocas máx.",       "100"),
    ("Batch size",        "2 048"),
    ("Monitor HPO",       "AUC-ROC (val)"),
    ("Diferenciable",     "✓  Sí — white-box OK"),
]
for i, (k, v) in enumerate(specs):
    ry = 1.38 + i * 0.44
    label(sl, k, 8.45, ry, 2.0, 0.40, size=9, bold=True, color=SUB_CLR, align=PP_ALIGN.LEFT)
    label(sl, v, 10.45, ry, 2.5, 0.40, size=9, color=WHITE, align=PP_ALIGN.LEFT)

note = ("MLP con 3 capas ocultas de 128 neuronas cada una. "
        "La BatchNorm en entrada estabiliza el entrenamiento con datos tabulares. "
        "La pérdida ponderada compensa el desbalance de clases.")
box(sl, 8.35, 6.08, 4.6, 1.1, fill=RGBColor(0x0A, 0x2A, 0x18), border=ACC_GREEN)
label(sl, note, 8.45, 6.1, 4.45, 1.05, size=8, color=SUB_CLR, italic=True, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# Slide 4 — LSTM-Attention
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "LSTM con Mecanismo de Atención",
            "Red recurrente con atención aditiva para capturar dependencias temporales", ACC_PURP)
footer(sl)

box(sl, 0.25, 0.95, 7.8, 6.3, border=ACC_PURP)
label(sl, "ARQUITECTURA  (hidden_dim=64, n_lstm_layers=2, dropout=0.3)", 0.35, 0.98, 7.5, 0.28,
      size=8, color=ACC_PURP, bold=True, align=PP_ALIGN.LEFT)

# Bloques del flujo
layers_lstm = [
    ("INPUT\n(batch, input_dim)",           SUB_CLR,  1.05, 1.6),
    ("Reshape\n(batch, seq_len, step_dim)", ACC_BLUE, 1.05, 2.1),
    ("LSTM × 2 capas\nhidden=64  dropout=0.3\n(batch, seq, 64)", ACC_PURP, 1.05, 3.05),
    ("Atención Aditiva\ntanh → Linear(64,1)\n→ Softmax → Σwᵢhᵢ", RGBColor(0xEC,0x4E,0x99), 1.05, 4.25),
    ("context vector\n(batch, 64)",          SUB_CLR,  1.05, 5.4),
    ("Dropout(0.3)\nLinear(64→2)",           ACC_PURP, 1.05, 5.95),
    ("OUTPUT  logits\n(batch, 2)",           SUB_CLR,  1.05, 6.75),
]

for i, (lbl, clr, bx_l, by_l) in enumerate(layers_lstm):
    bw_l = 2.5
    bh_l = 0.72 if "LSTM" in lbl or "Atención" in lbl else 0.50
    s = sl.shapes.add_shape(5, Inches(bx_l), Inches(by_l), Inches(bw_l), Inches(bh_l))
    s.fill.solid(); s.fill.fore_color.rgb = CARD
    s.line.color.rgb = clr; s.line.width = Pt(1.8 if clr == ACC_PURP else 1.2)
    label(sl, lbl, bx_l, by_l+0.02, bw_l, bh_l-0.04, size=8.5, bold=(clr != SUB_CLR), color=clr)
    if i < len(layers_lstm) - 1:
        next_by = layers_lstm[i+1][3]
        arrow(sl, bx_l + bw_l/2, by_l + bh_l, bx_l + bw_l/2, next_by, ACC_PURP)

# Detalle del bloque LSTM
box(sl, 3.8, 2.9, 4.1, 1.8, fill=RGBColor(0x17, 0x10, 0x3A), border=ACC_PURP)
label(sl, "CELDA LSTM", 3.9, 2.92, 3.9, 0.28, size=8, bold=True, color=ACC_PURP, align=PP_ALIGN.LEFT)
gates = [("Olvido (f)", "σ(Wf·[h,x]+bf)"),
         ("Entrada (i)", "σ(Wi·[h,x]+bi)"),
         ("Celda (c̃)",   "tanh(Wc·[h,x]+bc)"),
         ("Salida (o)",   "σ(Wo·[h,x]+bo)")]
for j, (g, formula) in enumerate(gates):
    gy = 3.25 + j * 0.35
    label(sl, g,       3.9,  gy, 1.3, 0.32, size=8, bold=True, color=ACC_PURP, align=PP_ALIGN.LEFT)
    label(sl, formula, 5.2,  gy, 2.6, 0.32, size=8, color=WHITE, align=PP_ALIGN.LEFT)

# Detalle atención
box(sl, 3.8, 4.8, 4.1, 1.4, fill=RGBColor(0x2A, 0x10, 0x2A), border=RGBColor(0xEC,0x4E,0x99))
label(sl, "ATENCIÓN ADITIVA", 3.9, 4.82, 3.9, 0.28, size=8, bold=True,
      color=RGBColor(0xEC,0x4E,0x99), align=PP_ALIGN.LEFT)
att_steps = [
    "1. score(hₜ) = tanh(W·hₜ + b)   →   e ∈ ℝ",
    "2. αₜ = softmax(e)               →   pesos de atención",
    "3. c  = Σ αₜ · hₜ               →   vector de contexto",
]
for j, step in enumerate(att_steps):
    label(sl, step, 3.9, 5.16 + j * 0.32, 3.9, 0.30, size=8, color=WHITE, align=PP_ALIGN.LEFT)

# Panel derecho
box(sl, 8.3, 0.95, 4.75, 6.3, border=BORDER)
label(sl, "ESPECIFICACIONES", 8.45, 0.98, 4.5, 0.28, size=8, color=ACC_PURP, bold=True, align=PP_ALIGN.LEFT)

specs = [
    ("Capas LSTM",        "n_lstm_layers (HPO: 1–4)"),
    ("Neuronas LSTM",     "hidden_dim (HPO: 32–256)"),
    ("Dropout",           "HPO: [0.0, 0.5]"),
    ("Atención",          "Aditiva  (Bahdanau-like)"),
    ("Activación att.",   "tanh + Softmax"),
    ("Seq. length",       "1 (dato tabular como seq.)"),
    ("Activación salida", "Softmax (en pérdida)"),
    ("Pérdida",           "Cross-Entropy ponderada"),
    ("Optimizador",       "Adam"),
    ("lr",                "HPO: [1e-5, 1e-2]"),
    ("weight_decay",      "HPO: [1e-6, 1e-2]"),
    ("Épocas máx.",       "100"),
    ("Batch size",        "1 024"),
    ("Diferenciable",     "✓  Sí — white-box OK"),
]
for i, (k, v) in enumerate(specs):
    ry = 1.38 + i * 0.40
    label(sl, k, 8.45, ry, 1.9, 0.38, size=8.5, bold=True, color=SUB_CLR, align=PP_ALIGN.LEFT)
    label(sl, v, 10.35, ry, 2.6, 0.38, size=8.5, color=WHITE, align=PP_ALIGN.LEFT)

note = ("LSTM con atención aditiva. Al tratar cada transacción como "
        "una secuencia de longitud 1, el mecanismo de atención pondera "
        "la importancia de cada paso oculto del LSTM.")
box(sl, 8.35, 6.08, 4.6, 1.1, fill=RGBColor(0x1A, 0x0D, 0x3A), border=ACC_PURP)
label(sl, note, 8.45, 6.1, 4.45, 1.05, size=8, color=SUB_CLR, italic=True, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# Slide 5 — XGBoost
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "XGBoost — Gradient Boosting de Árboles",
            "Ensamble de árboles de decisión entrenados con boosting por gradiente", ACC_ORG)
footer(sl)

box(sl, 0.25, 0.95, 7.8, 6.3, border=ACC_ORG)
label(sl, "ARQUITECTURA  (n_estimators=100, max_depth=6)", 0.35, 0.98, 7.5, 0.28,
      size=8, color=ACC_ORG, bold=True, align=PP_ALIGN.LEFT)

# Flujo general
flow = [
    ("INPUT\nvector de features x", SUB_CLR, 0.8, 1.15, 2.4, 0.52),
    ("Árbol 1\n(residuo inicial)",  ACC_ORG, 0.5, 2.1,  1.7, 0.65),
    ("Árbol 2\n(corrige errores)",  ACC_ORG, 2.4, 2.1,  1.7, 0.65),
    ("···",                          SUB_CLR, 4.2, 2.25, 0.5, 0.40),
    ("Árbol K\n(K=100)",            ACC_ORG, 4.8, 2.1,  1.7, 0.65),
    ("Σ predicciones\nponderadas",  RED_CLR, 1.5, 3.25, 2.2, 0.55),
    ("Sigmoid / Softprob\nP(fraude)", ACC_ORG, 1.5, 4.25, 2.2, 0.55),
    ("OUTPUT\nP(legítimo) / P(fraude)", SUB_CLR, 1.5, 5.2, 2.2, 0.52),
]
for lbl, clr, fx, fy, fw, fh in flow:
    if lbl == "···":
        label(sl, lbl, fx, fy, fw, fh, size=16, bold=True, color=SUB_CLR)
        continue
    s = sl.shapes.add_shape(5, Inches(fx), Inches(fy), Inches(fw), Inches(fh))
    s.fill.solid(); s.fill.fore_color.rgb = CARD
    s.line.color.rgb = clr; s.line.width = Pt(1.5)
    label(sl, lbl, fx, fy+0.04, fw, fh-0.08, size=9, bold=(clr != SUB_CLR), color=clr)

# Flechas del flujo
arrow(sl, 2.0, 1.67, 1.35, 2.1, ACC_ORG)
arrow(sl, 2.0, 1.67, 3.25, 2.1, ACC_ORG)
arrow(sl, 2.0, 1.67, 5.65, 2.1, ACC_ORG)
arrow(sl, 1.35, 2.75, 2.6, 3.25, ACC_ORG)
arrow(sl, 3.25, 2.75, 2.6, 3.25, ACC_ORG)
arrow(sl, 5.65, 2.75, 2.6, 3.25, ACC_ORG)
arrow(sl, 2.6,  3.80, 2.6, 4.25, ACC_ORG)
arrow(sl, 2.6,  4.80, 2.6, 5.20, ACC_ORG)

# Detalle de un árbol
box(sl, 4.0, 3.2, 3.9, 3.8, fill=RGBColor(0x2A, 0x1A, 0x08), border=ACC_ORG)
label(sl, "DETALLE: ÁRBOL DE DECISIÓN", 4.1, 3.22, 3.7, 0.28,
      size=8, bold=True, color=ACC_ORG, align=PP_ALIGN.LEFT)

# Árbol simple dibujado
tree_root_x, tree_root_y = 5.95, 3.75
circle(sl, tree_root_x, tree_root_y, 0.28, RGBColor(0x3A,0x2A,0x00), ACC_ORG)
label(sl, "x[j]>θ", tree_root_x-0.28, tree_root_y-0.2, 0.56, 0.4, size=7.5, color=ACC_ORG)

lx1, ly1 = 5.25, 4.45
lx2, ly2 = 6.65, 4.45
circle(sl, lx1, ly1, 0.24, CARD, ACC_ORG)
circle(sl, lx2, ly2, 0.24, CARD, ACC_ORG)
label(sl, "hoja\n0.3", lx1-0.24, ly1-0.2, 0.48, 0.4, size=7, color=WHITE)
label(sl, "hoja\n0.8", lx2-0.24, ly2-0.2, 0.48, 0.4, size=7, color=WHITE)
arrow(sl, tree_root_x-0.28, tree_root_y+0.1, lx1+0.06, ly1-0.24, ACC_ORG)
arrow(sl, tree_root_x+0.28, tree_root_y+0.1, lx2-0.06, ly2-0.24, ACC_ORG)
label(sl, "SÍ", 5.3, 4.08, 0.35, 0.24, size=7, color=ACC_GREEN, bold=True)
label(sl, "NO", 6.5, 4.08, 0.35, 0.24, size=7, color=RED_CLR, bold=True)
label(sl, "max_depth=6\n→ hasta 64 hojas", 4.1, 4.85, 3.7, 0.40,
      size=8, color=SUB_CLR, italic=True)

# Por qué NO es diferenciable
box(sl, 4.05, 5.35, 3.8, 1.75, fill=RGBColor(0x2A, 0x08, 0x08), border=RED_CLR)
label(sl, "⚠  NO DIFERENCIABLE", 4.15, 5.37, 3.6, 0.28, size=8.5, bold=True, color=RED_CLR, align=PP_ALIGN.LEFT)
label(sl,
      "Las decisiones x[j] > θ son funciones\n"
      "escalón (discontinuas). No existe ∂f/∂x.\n"
      "→ Incompatible con ataques white-box (CaFA).\n"
      "→ Compatible con ataques black-box\n"
      "   (HopSkipJump, BoundaryAttack, SquareAttack).",
      4.15, 5.68, 3.65, 1.35, size=8, color=WHITE, align=PP_ALIGN.LEFT)

# Panel derecho
box(sl, 8.3, 0.95, 4.75, 6.3, border=BORDER)
label(sl, "ESPECIFICACIONES", 8.45, 0.98, 4.5, 0.28, size=8, color=ACC_ORG, bold=True, align=PP_ALIGN.LEFT)

specs = [
    ("Tipo",              "XGBClassifier (xgboost)"),
    ("Estimadores",       "n_estimators: HPO 100–600"),
    ("Prof. max.",        "max_depth: HPO 3–10"),
    ("Learning rate",     "HPO: [0.01, 0.3]"),
    ("Objetivo",          "binary:logistic"),
    ("Métrica eval.",     "AUC-PR (aucpr)"),
    ("Desbalance",        "scale_pos_weight\n= class_weight_minority"),
    ("min_child_weight",  "HPO: [1, 20]"),
    ("subsample",         "HPO: [0.5, 1.0]"),
    ("colsample_bytree",  "HPO: [0.5, 1.0]"),
    ("HPO trials",        "30 (Optuna TPE, seed=42)"),
    ("Diferenciable",     "✗  No — solo black-box"),
]
for i, (k, v) in enumerate(specs):
    ry = 1.38 + i * 0.44
    label(sl, k, 8.45, ry, 2.0, 0.40, size=9, bold=True, color=SUB_CLR, align=PP_ALIGN.LEFT)
    label(sl, v, 10.45, ry, 2.5, 0.40, size=8.5, color=WHITE, align=PP_ALIGN.LEFT)

note = ("XGBoost es el modelo más robusto en escenarios con desbalance "
        "extremo gracias a scale_pos_weight. Al no ser diferenciable, "
        "solo puede ser evaluado con ataques de caja negra.")
box(sl, 8.35, 6.08, 4.6, 1.1, fill=RGBColor(0x2A, 0x18, 0x08), border=ACC_ORG)
label(sl, note, 8.45, 6.1, 4.45, 1.05, size=8, color=SUB_CLR, italic=True, align=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════════════════════════════
# Slide 6 — Comparativa
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
slide_title(sl, "Comparativa de Modelos",
            "Resumen de capacidades, compatibilidad con ataques y características clave", ACC_BLUE)
footer(sl)

box(sl, 0.25, 0.95, 12.85, 6.3, border=BORDER)

headers = ["Característica", "Log. Reg.", "MLP", "LSTM-Att.", "XGBoost"]
hcolors = [SUB_CLR, ACC_BLUE, ACC_GREEN, ACC_PURP, ACC_ORG]
col_w   = [2.8, 2.3, 2.3, 2.3, 2.3]
col_x   = [0.35]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

# Header row
for j, (h, clr, cx, cw) in enumerate(zip(headers, hcolors, col_x, col_w)):
    s = sl.shapes.add_shape(1, Inches(cx), Inches(1.12), Inches(cw-0.06), Inches(0.42))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x17, 0x24, 0x38) if j == 0 else CARD
    s.line.color.rgb = clr; s.line.width = Pt(1.2)
    label(sl, h, cx, 1.16, cw-0.06, 0.34, size=10, bold=True, color=clr)

rows = [
    ("Tipo de modelo",    "Lineal",          "Red neuronal\nfeed-forward", "Red recurrente\n+ Atención", "Ensamble\nárboles"),
    ("Capas / Profund.",  "1 capa (lineal)", "3 ocultas × 128\n(configurable)", "2 LSTM × 64\n(configurable)", "100 árboles\nprof. 6"),
    ("Activación",        "Softmax (salida)", "ReLU + Softmax", "tanh + Softmax", "Funciones\nescalón"),
    ("Diferenciable",     "✓  Sí",           "✓  Sí",          "✓  Sí",          "✗  No"),
    ("CaFA (white-box)",  "✓",               "✓",              "✓",              "✗"),
    ("HopSkipJump",       "✓",               "✓",              "✓",              "✓"),
    ("BoundaryAttack",    "✓",               "✓",              "✓",              "✓"),
    ("SquareAttack",      "✓",               "✓",              "✓",              "✓"),
    ("Desbalance",        "class_weight",    "class_weight",   "class_weight",   "scale_pos_weight"),
    ("SMOTE requerido",   "Opcional",        "Necesario",      "Necesario",      "Opcional"),
    ("Optimizador",       "Adam",            "Adam",           "Adam",           "XGBoost (HPO)"),
    ("HPO tool",          "Optuna",          "Optuna",         "Optuna",         "Optuna TPE"),
]

row_colors_alt = [RGBColor(0x17, 0x22, 0x35), CARD]
for i, row in enumerate(rows):
    ry = 1.62 + i * 0.40
    for j, (val, cx, cw) in enumerate(zip(row, col_x, col_w)):
        s = sl.shapes.add_shape(1, Inches(cx), Inches(ry), Inches(cw-0.06), Inches(0.38))
        s.fill.solid()
        s.fill.fore_color.rgb = row_colors_alt[i % 2]
        s.line.color.rgb = BORDER; s.line.width = Pt(0.5)

        # Color especial para checkmarks/X
        clr = WHITE
        if val == "✓  Sí" or val == "✓": clr = ACC_GREEN
        elif val == "✗  No" or val == "✗": clr = RED_CLR
        elif j == 0: clr = SUB_CLR
        elif j == 1: clr = ACC_BLUE
        elif j == 2: clr = ACC_GREEN
        elif j == 3: clr = ACC_PURP
        elif j == 4: clr = ACC_ORG

        label(sl, val, cx, ry+0.02, cw-0.06, 0.35, size=8.5,
              color=clr, bold=(j==0))


# ════════════════════════════════════════════════════════════════════════════
# Guardar
# ════════════════════════════════════════════════════════════════════════════
OUT = "/home/omar/Desktop/TESIS1/adversarial-attack-cafa/presentacion/arquitecturas_modelos.pptx"
prs.save(OUT)
print(f"Guardado: {OUT}")
